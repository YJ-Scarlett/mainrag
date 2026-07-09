import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from docx import Document
from fastapi import HTTPException
from pypdf import PdfReader
from pptx import Presentation

from core.config import settings


SUPPORTED_EXTENSIONS = {".doc", ".docx", ".ppt", ".pptx", ".pdf"}


def _extract_pdf(path: Path) -> str:
    try:
        return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
    except Exception as exc:
        raise HTTPException(400, "PDF 解析失败，请确认文件未损坏且不是加密文件。") from exc


def _extract_docx(path: Path) -> str:
    try:
        document = Document(path)
        lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if text:
                    lines.append(text)
        return "\n".join(lines)
    except Exception as exc:
        raise HTTPException(400, "DOCX 解析失败，请确认文件未损坏。") from exc


def _extract_pptx(path: Path) -> str:
    try:
        presentation = Presentation(path)
        slides = []
        for number, slide in enumerate(presentation.slides, 1):
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if text:
                            lines.append(text)
            if lines:
                slides.append(f"第 {number} 页\n" + "\n".join(lines))
        return "\n\n".join(slides)
    except Exception as exc:
        raise HTTPException(400, "PPTX 解析失败，请确认文件未损坏。") from exc


def _convert_legacy_office(source: Path, target_extension: str) -> Path:
    """使用 Windows 本机 Microsoft Office 将 DOC/PPT 转为 DOCX/PPTX。"""
    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:
        raise HTTPException(500, "缺少 pywin32，无法解析旧版 DOC/PPT 文件。") from exc

    work_dir = Path(tempfile.mkdtemp(prefix="zhiwen-office-"))
    target = work_dir / f"converted{target_extension}"
    pythoncom.CoInitialize()
    application = None
    try:
        if source.suffix.lower() == ".doc":
            application = win32com.client.DispatchEx("Word.Application")
            application.Visible = False
            application.DisplayAlerts = 0
            document = application.Documents.Open(str(source.resolve()), ReadOnly=True)
            try:
                document.SaveAs2(str(target), FileFormat=16)  # wdFormatDocumentDefault (.docx)
            finally:
                document.Close(False)
        else:
            application = win32com.client.DispatchEx("PowerPoint.Application")
            presentation = application.Presentations.Open(str(source.resolve()), ReadOnly=True, Untitled=False, WithWindow=False)
            try:
                presentation.SaveAs(str(target), 24)  # ppSaveAsOpenXMLPresentation (.pptx)
            finally:
                presentation.Close()
        return target
    except Exception as exc:
        shutil.rmtree(work_dir, ignore_errors=True)
        raise HTTPException(
            400,
            "旧版 DOC/PPT 解析失败，请确认本机已安装 Microsoft Office，或另存为 DOCX/PPTX 后再上传。",
        ) from exc
    finally:
        if application is not None:
            application.Quit()
        pythoncom.CoUninitialize()


def create_pdf_preview(source: Path, output: Path) -> Path:
    """生成保留原始版式的 PDF 预览文件。"""
    suffix = source.suffix.lower()
    if suffix == ".pdf":
        shutil.copy2(source, output)
        return output

    try:
        completed = subprocess.run(
            [
                sys.executable,
                "-m",
                "services.preview_worker",
                str(source.resolve()),
                str(output.resolve()),
            ],
            cwd=str(Path(__file__).resolve().parents[1]),
            capture_output=True,
            text=True,
            timeout=settings.preview_timeout_seconds,
        )
    except subprocess.TimeoutExpired as exc:
        raise HTTPException(408, "Office 生成预览超时，已跳过原版式预览。") from exc

    if completed.returncode != 0 or not output.exists():
        detail = (completed.stderr or completed.stdout or "").strip()
        raise HTTPException(
            400,
            f"Office 生成预览失败，已跳过原版式预览。{detail[:180]}",
        )
    return output


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise HTTPException(400, "仅支持 DOC、DOCX、PPT、PPTX 和 PDF 文件。")

    converted = None
    try:
        if suffix == ".pdf":
            return _extract_pdf(path)
        if suffix == ".docx":
            return _extract_docx(path)
        if suffix == ".pptx":
            return _extract_pptx(path)
        converted = _convert_legacy_office(path, ".docx" if suffix == ".doc" else ".pptx")
        return _extract_docx(converted) if suffix == ".doc" else _extract_pptx(converted)
    finally:
        if converted is not None:
            shutil.rmtree(converted.parent, ignore_errors=True)
